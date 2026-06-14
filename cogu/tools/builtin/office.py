import csv
import io
import os
from pathlib import Path
from typing import Any

from cogu.tools.base import FunctionTool, ToolRegistry, ToolCapability


def _pdf_read(path: str, pages: str = "") -> str:
    p = Path(path)
    if not p.exists():
        return f"Error: file not found: {path}"
    try:
        from pypdf import PdfReader
    except ImportError:
        return "Error: pypdf not installed. Run: pip install pypdf"
    reader = PdfReader(str(p))
    total = len(reader.pages)
    page_nums = []
    if pages:
        for part in pages.split(","):
            part = part.strip()
            if "-" in part:
                a, b = part.split("-", 1)
                page_nums.extend(range(int(a) - 1, int(b)))
            else:
                page_nums.append(int(part) - 1)
    else:
        page_nums = list(range(min(total, 50)))
    text_parts = []
    for i in page_nums:
        if 0 <= i < total:
            text_parts.append(f"--- Page {i + 1} ---\n{reader.pages[i].extract_text()}")
    result = "\n\n".join(text_parts)
    if len(page_nums) < total:
        result += f"\n\n[Showing pages {page_nums[0] + 1}-{page_nums[-1] + 1} of {total}]"
    return result


def _pdf_merge(output: str, inputs: str) -> str:
    try:
        from pypdf import PdfWriter, PdfReader
    except ImportError:
        return "Error: pypdf not installed. Run: pip install pypdf"
    writer = PdfWriter()
    for inp in inputs.split(","):
        inp = inp.strip()
        if not Path(inp).exists():
            return f"Error: file not found: {inp}"
        reader = PdfReader(inp)
        for page in reader.pages:
            writer.add_page(page)
    Path(output).parent.mkdir(parents=True, exist_ok=True)
    writer.write(output)
    return f"Merged {len(inputs.split(','))} files -> {output}"


def _pdf_split(path: str, output_dir: str, pages_per: int = 1) -> str:
    try:
        from pypdf import PdfReader, PdfWriter
    except ImportError:
        return "Error: pypdf not installed. Run: pip install pypdf"
    reader = PdfReader(path)
    total = len(reader.pages)
    out = Path(output_dir)
    out.mkdir(parents=True, exist_ok=True)
    for i in range(0, total, pages_per):
        writer = PdfWriter()
        for j in range(i, min(i + pages_per, total)):
            writer.add_page(reader.pages[j])
        fname = out / f"split_{i + 1}_{min(i + pages_per, total)}.pdf"
        writer.write(str(fname))
    return f"Split {total} pages into {out}"


def _docx_read(path: str) -> str:
    p = Path(path)
    if not p.exists():
        return f"Error: file not found: {path}"
    try:
        from docx import Document
    except ImportError:
        return "Error: python-docx not installed. Run: pip install python-docx"
    doc = Document(str(p))
    text = "\n".join(para.text for para in doc.paragraphs if para.text.strip())
    if not text:
        return "[document contains no text paragraphs]"
    if len(text) > 10000:
        text = text[:10000] + "\n\n[truncated]"
    return text


def _xlsx_read(path: str, sheet: str = "") -> str:
    p = Path(path)
    if not p.exists():
        return f"Error: file not found: {path}"
    try:
        import openpyxl
    except ImportError:
        return "Error: openpyxl not installed. Run: pip install openpyxl"
    wb = openpyxl.load_workbook(str(p), data_only=True)
    ws = wb[sheet] if sheet else wb.active
    output = io.StringIO()
    writer = csv.writer(output)
    for row in ws.iter_rows(values_only=True):
        writer.writerow([str(c) if c is not None else "" for c in row])
    result = output.getvalue()
    if len(result) > 10000:
        result = result[:10000] + "\n\n[truncated]"
    return result


def _docx_write(path: str, content: str, title: str = "") -> str:
    p = Path(path)
    try:
        from docx import Document
        from docx.shared import Pt
    except ImportError:
        return "Error: python-docx not installed. Run: pip install python-docx"
    doc = Document()
    if title:
        heading = doc.add_heading(title, level=1)
        heading.alignment = 1
    lines = content.split("\n")
    for line in lines:
        line = line.rstrip()
        if not line:
            doc.add_paragraph()
            continue
        if line.startswith("# "):
            doc.add_heading(line[2:], level=1)
        elif line.startswith("## "):
            doc.add_heading(line[3:], level=2)
        elif line.startswith("### "):
            doc.add_heading(line[4:], level=3)
        elif line.startswith("- ") or line.startswith("* "):
            doc.add_paragraph(line[2:], style="List Bullet")
        elif line.startswith("1. ") or line.startswith("2. ") or line.startswith("3. "):
            doc.add_paragraph(line[3:], style="List Number")
        else:
            doc.add_paragraph(line)
    p.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(p))
    return f"Created Word document: {path}"


def _xlsx_write(path: str, data: str, sheet_name: str = "Sheet1") -> str:
    p = Path(path)
    try:
        import openpyxl
    except ImportError:
        return "Error: openpyxl not installed. Run: pip install openpyxl"
    wb = openpyxl.Workbook()
    if sheet_name and sheet_name != "Sheet1":
        wb.remove(wb.active)
        wb.create_sheet(sheet_name)
    ws = wb.active
    if sheet_name:
        ws.title = sheet_name
    reader = csv.reader(io.StringIO(data))
    for row_idx, row in enumerate(reader, start=1):
        for col_idx, value in enumerate(row, start=1):
            ws.cell(row=row_idx, column=col_idx, value=value)
    p.parent.mkdir(parents=True, exist_ok=True)
    wb.save(str(p))
    return f"Created Excel spreadsheet: {path}"


def _pptx_write(path: str, slides: str, title: str = "Presentation") -> str:
    p = Path(path)
    try:
        from pptx import Presentation
    except ImportError:
        return "Error: python-pptx not installed. Run: pip install python-pptx"
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = "Generated by COGU Agent"
    slide_texts = slides.split("---slide---")
    for slide_text in slide_texts[1:]:
        slide_text = slide_text.strip()
        if not slide_text:
            continue
        lines = slide_text.split("\n", 1)
        slide_title = lines[0].strip() if lines else ""
        slide_content = lines[1].strip() if len(lines) > 1 else ""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        if slide_title:
            slide.shapes.title.text = slide_title
        if slide_content:
            tf = slide.placeholders[1].text_frame
            tf.text = slide_content
    p.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(p))
    return f"Created PowerPoint presentation: {path}"


def _markdown_to_docx(markdown_path: str, docx_path: str) -> str:
    md_p = Path(markdown_path)
    if not md_p.exists():
        return f"Error: file not found: {markdown_path}"
    try:
        from docx import Document
        from docx.shared import Pt
    except ImportError:
        return "Error: python-docx not installed. Run: pip install python-docx"
    content = md_p.read_text(encoding="utf-8")
    return _docx_write(docx_path, content)


def _report_docx(path: str, title: str = "", sections: str = "") -> str:
    lines = []
    if title:
        lines.append(f"# {title}")
    if sections:
        lines.append(sections)
    content = "\n\n".join(lines)
    return _docx_write(path, content, title)


def register_office_tools(registry: ToolRegistry):
    registry.register(FunctionTool(_pdf_read, name="pdf_read", description="Read text from a PDF file. Specify pages like '1-5' or '1,3,5'.").with_capability(ToolCapability.READ_ONLY).mark_concurrency_safe().with_group("office"))
    registry.register(FunctionTool(_pdf_merge, name="pdf_merge", description="Merge multiple PDF files into one. Inputs as comma-separated paths.").with_capability(ToolCapability.WRITES_FILES).with_group("office"))
    registry.register(FunctionTool(_pdf_split, name="pdf_split", description="Split a PDF into multiple files by page count.").with_capability(ToolCapability.WRITES_FILES).with_group("office"))
    registry.register(FunctionTool(_docx_read, name="docx_read", description="Read text content from a .docx Word document.").with_capability(ToolCapability.READ_ONLY).mark_concurrency_safe().with_group("office"))
    registry.register(FunctionTool(_xlsx_read, name="xlsx_read", description="Read data from an Excel spreadsheet as CSV. Specify sheet name optionally.").with_capability(ToolCapability.READ_ONLY).mark_concurrency_safe().with_group("office"))
    registry.register(FunctionTool(_docx_write, name="docx_write", description="Create a Word document from content. Supports basic markdown: # heading, ## subheading, - list, 1. numbered list.").with_capability(ToolCapability.WRITES_FILES).with_group("office"))
    registry.register(FunctionTool(_xlsx_write, name="xlsx_write", description="Create an Excel spreadsheet from CSV data. First row becomes headers.").with_capability(ToolCapability.WRITES_FILES).with_group("office"))
    registry.register(FunctionTool(_pptx_write, name="pptx_write", description="Create a PowerPoint presentation. Format: title slide then ---slide--- to separate slides, each with title\\ncontent.").with_capability(ToolCapability.WRITES_FILES).with_group("office"))
    registry.register(FunctionTool(_markdown_to_docx, name="markdown_to_docx", description="Convert a markdown file to a Word document.").with_capability(ToolCapability.WRITES_FILES).with_group("office"))
    registry.register(FunctionTool(_report_docx, name="report_docx", description="Create a structured report document. Title and sections (each section starts with ##).").with_capability(ToolCapability.WRITES_FILES).with_group("office"))
