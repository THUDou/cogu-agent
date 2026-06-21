#!/usr/bin/env python3
"""
Convert HTML slide files to a PPTX presentation.
Uses Playwright to render HTML as screenshots, then python-pptx to assemble.
"""

import os
import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
from playwright.sync_api import sync_playwright
import argparse


def natural_sort_key(filename):
    """Natural sort so page_1 comes before page_10."""
    numbers = re.findall(r'\d+', filename)
    return int(numbers[0]) if numbers else 0


SCALE_SCRIPT = """
() => {
    const container = document.querySelector('.slide-container') ||
                    document.querySelector('#slide-container') ||
                    document.querySelector('[class*="slide"]') ||
                    document.querySelector('[id*="slide"]');

    if (container) {
        const containerWidth = parseFloat(getComputedStyle(container).width);
        const containerHeight = parseFloat(getComputedStyle(container).height);
        const viewportWidth = window.innerWidth;
        const viewportHeight = window.innerHeight;

        const scaleX = viewportWidth / containerWidth;
        const scaleY = viewportHeight / containerHeight;
        const scale = Math.min(scaleX, scaleY);

        container.style.transform = `scale(${scale})`;
        container.style.transformOrigin = 'center center';

        document.body.style.display = 'flex';
        document.body.style.justifyContent = 'center';
        document.body.style.alignItems = 'center';
        document.body.style.width = '100vw';
        document.body.style.height = '100vh';
        document.body.style.overflow = 'hidden';

        return true;
    }
    return false;
}
"""


def html_to_image(browser, html_path, output_image_path, width=1920, height=1080, wait_time=3000):
    """Render an HTML file to a PNG screenshot using Playwright."""
    page = browser.new_page(viewport={'width': width, 'height': height})
    page.goto(f'file://{html_path}')
    page.wait_for_load_state('networkidle')
    page.evaluate(SCALE_SCRIPT)
    page.wait_for_timeout(wait_time)
    page.screenshot(path=output_image_path, full_page=False)
    page.close()
    print(f"  Screenshot: {output_image_path}")


def create_ppt_from_images(image_dir, output_ppt_path, width=1920, height=1080):
    """Assemble PNG screenshots into a PPTX presentation."""
    prs = Presentation()
    prs.slide_width = Inches(width / 96)
    prs.slide_height = Inches(height / 96)

    image_files = sorted(
        [f for f in os.listdir(image_dir) if f.endswith('.png')],
        key=natural_sort_key
    )

    print(f"  Found {len(image_files)} images")

    for image_file in image_files:
        image_path = os.path.join(image_dir, image_file)
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        slide.shapes.add_picture(
            image_path, 0, 0,
            width=prs.slide_width,
            height=prs.slide_height
        )
        print(f"  Added slide: {image_file}")

    prs.save(output_ppt_path)
    print(f"  Saved: {output_ppt_path}")


def convert_html_to_ppt(html_dir, output_ppt_path, temp_dir=None, wait_time=3000, width=1920, height=1080):
    """Main: convert a directory of HTML files to a PPTX presentation."""
    html_dir = Path(html_dir).resolve()

    if temp_dir is None:
        temp_dir = html_dir / 'temp_images'
    else:
        temp_dir = Path(temp_dir)

    temp_dir.mkdir(exist_ok=True)

    html_files = sorted(
        [f for f in os.listdir(html_dir) if f.endswith('.html') and f.startswith('page_')],
        key=natural_sort_key
    )

    if not html_files:
        print("Error: No HTML files found (page_*.html)")
        return

    print(f"Converting {len(html_files)} HTML files to PPTX...")

    # Step 1: Render HTML to images
    print("Step 1: Rendering HTML to images...")
    with sync_playwright() as p:
        browser = p.chromium.launch()
        for html_file in html_files:
            html_path = html_dir / html_file
            image_name = html_file.replace('.html', '.png')
            image_path = temp_dir / image_name
            html_to_image(browser, str(html_path), str(image_path), width=width, height=height, wait_time=wait_time)
        browser.close()

    # Step 2: Assemble PPTX
    print("Step 2: Creating PPTX...")
    create_ppt_from_images(str(temp_dir), output_ppt_path, width=width, height=height)

    print(f"\nDone! Output: {output_ppt_path}")


def main():
    parser = argparse.ArgumentParser(description='Convert HTML slide files to PPTX')
    parser.add_argument('html_dir', help='Directory containing page_*.html files')
    parser.add_argument('-o', '--output', help='Output PPTX path (default: output.pptx)', default='output.pptx')
    parser.add_argument('-t', '--temp-dir', help='Temp directory for screenshots')
    parser.add_argument('-w', '--wait', type=int, default=3000, help='Wait time in ms before screenshot (default: 3000)')
    parser.add_argument('--width', type=int, default=1920, help='Viewport width (default: 1920)')
    parser.add_argument('--height', type=int, default=1080, help='Viewport height (default: 1080)')

    args = parser.parse_args()
    convert_html_to_ppt(args.html_dir, args.output, args.temp_dir, args.wait, args.width, args.height)


if __name__ == '__main__':
    main()
