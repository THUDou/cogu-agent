
import argparse
import shutil
import sys
from copy import deepcopy
from pathlib import Path

import six
from pptx import Presentation


def main():
    parser = argparse.ArgumentParser(
        description="Rearrange PowerPoint slides based on a sequence of indices.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python rearrange.py template.pptx output.pptx 0,34,34,50,52
    Creates output.pptx using slides 0, 34 (twice), 50, and 52 from template.pptx

  python rearrange.py template.pptx output.pptx 5,3,1,2,4
    Creates output.pptx with slides reordered as specified

Note: Slide indices are 0-based (first slide is 0, second is 1, etc.)
    source = pres.slides[index]

    new_slide = pres.slides.add_slide(source.slide_layout)

    image_rels = {}
    for rel_id, rel in six.iteritems(source.part.rels):
        if "image" in rel.reltype or "media" in rel.reltype:
            image_rels[rel_id] = rel

    for shape in new_slide.shapes:
        sp = shape.element
        sp.getparent().remove(sp)

    for shape in source.shapes:
        el = shape.element
        new_el = deepcopy(el)
        new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")

        blips = new_el.xpath(".//a:blip[@r:embed]")
        for blip in blips:
            old_rId = blip.get(
                "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
            )
            if old_rId in image_rels:
                old_rel = image_rels[old_rId]
                new_rId = new_slide.part.rels.get_or_add(
                    old_rel.reltype, old_rel._target
                )
                blip.set(
                    "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed",
                    new_rId,
                )

    for rel_id, rel in image_rels.items():
        try:
            new_slide.part.rels.get_or_add(rel.reltype, rel._target)
        except Exception:
            pass  # Relationship might already exist

    return new_slide


def delete_slide(pres, index):
    rId = pres.slides._sldIdLst[index].rId
    pres.part.drop_rel(rId)
    del pres.slides._sldIdLst[index]


def reorder_slides(pres, slide_index, target_index):
    slides = pres.slides._sldIdLst

    slide_element = slides[slide_index]
    slides.remove(slide_element)

    slides.insert(target_index, slide_element)


def rearrange_presentation(template_path, output_path, slide_sequence):
    if template_path != output_path:
        shutil.copy2(template_path, output_path)
        prs = Presentation(output_path)
    else:
        prs = Presentation(template_path)

    total_slides = len(prs.slides)

    for idx in slide_sequence:
        if idx < 0 or idx >= total_slides:
            raise ValueError(f"Slide index {idx} out of range (0-{total_slides - 1})")

    slide_map = []  # List of actual slide indices for final presentation
    duplicated = {}  # Track duplicates: original_idx -> [duplicate_indices]

    print(f"Processing {len(slide_sequence)} slides from template...")
    for i, template_idx in enumerate(slide_sequence):
        if template_idx in duplicated and duplicated[template_idx]:
            slide_map.append(duplicated[template_idx].pop(0))
            print(f"  [{i}] Using duplicate of slide {template_idx}")
        elif slide_sequence.count(template_idx) > 1 and template_idx not in duplicated:
            slide_map.append(template_idx)
            duplicates = []
            count = slide_sequence.count(template_idx) - 1
            print(
                f"  [{i}] Using original slide {template_idx}, creating {count} duplicate(s)"
            )
            for _ in range(count):
                duplicate_slide(prs, template_idx)
                duplicates.append(len(prs.slides) - 1)
            duplicated[template_idx] = duplicates
        else:
            slide_map.append(template_idx)
            print(f"  [{i}] Using original slide {template_idx}")

    slides_to_keep = set(slide_map)
    print(f"\nDeleting {len(prs.slides) - len(slides_to_keep)} unused slides...")
    for i in range(len(prs.slides) - 1, -1, -1):
        if i not in slides_to_keep:
            delete_slide(prs, i)
            slide_map = [idx - 1 if idx > i else idx for idx in slide_map]

    print(f"Reordering {len(slide_map)} slides to final sequence...")
    for target_pos in range(len(slide_map)):
        current_pos = slide_map[target_pos]
        if current_pos != target_pos:
            reorder_slides(prs, current_pos, target_pos)
            for i in range(len(slide_map)):
                if slide_map[i] > current_pos and slide_map[i] <= target_pos:
                    slide_map[i] -= 1
                elif slide_map[i] < current_pos and slide_map[i] >= target_pos:
                    slide_map[i] += 1
            slide_map[target_pos] = target_pos

    prs.save(output_path)
    print(f"\nSaved rearranged presentation to: {output_path}")
    print(f"Final presentation has {len(prs.slides)} slides")


if __name__ == "__main__":
    main()
