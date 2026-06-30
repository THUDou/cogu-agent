import sys
import os
import argparse
import logging
from pptx import Presentation


def split_pptx(input_path, output_folder, max_slides=0):
    if not os.path.exists(input_path):
        raise FileNotFoundError(f"输入文件不存在: {input_path}")
    os.makedirs(output_folder, exist_ok=True)

    prs_probe = Presentation(input_path)
    total_slides = len(prs_probe.slides)
    count = total_slides if max_slides <= 0 else min(total_slides, max_slides)
    logging.info("检测到共 %d 页，将拆分前 %d 页...", total_slides, count)

    for i in range(count):
        prs = Presentation(input_path)

        for j in range(total_slides - 1, i, -1):
            slide_id = prs.slides._sldIdLst[j].rId
            prs.part.drop_rel(slide_id)
            del prs.slides._sldIdLst[j]

        for j in range(i - 1, -1, -1):
            slide_id = prs.slides._sldIdLst[j].rId
            prs.part.drop_rel(slide_id)
            del prs.slides._sldIdLst[j]

        save_path = os.path.join(output_folder, f"slide_{i+1:03d}.pptx")
        prs.save(save_path)
        logging.info("已生成: %s", save_path)

    return count


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument('--input', required=True, help='PPTX 文件路径')
    parser.add_argument('--output', required=True, help='输出目录')
    parser.add_argument('--max-slides', type=int, default=10, dest='max_slides',
                        help='最多拆分页数（默认 10）')
    args = parser.parse_args()

    if not os.path.exists(args.input):
        logging.error("错误: 文件不存在: %s", args.input)
        sys.exit(1)

    try:
        count = split_pptx(args.input, args.output, args.max_slides)
        logging.info("拆分完成，共生成 %d 个单页文件", count)
    except Exception as e:
        logging.error("拆分失败: %s", e)
        sys.exit(1)


if __name__ == '__main__':
    main()
