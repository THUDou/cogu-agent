#!/usr/bin/env python3
"""
PPTX 结构提取脚本
从 PowerPoint 文件中提取主题、颜色、字体、布局等结构化数据
"""

import sys
import os
import json
import zipfile
import re
from pathlib import Path
from typing import Dict, List, Any, Optional
from collections import defaultdict, Counter

EMU_PER_INCH = 914400
EMU_PER_CM = 360000
EMU_PER_PT = 12700
NS       = 'http://schemas.openxmlformats.org/drawingml/2006/main'
NS_P_MAIN = 'http://schemas.openxmlformats.org/presentationml/2006/main'

# MSO_THEME_COLOR enum name → theme XML key（用于 schemeClr 解析）
_TC_KEY_MAP = {
    'BACKGROUND_1': 'lt1', 'LIGHT_1': 'lt1',
    'BACKGROUND_2': 'lt2', 'LIGHT_2': 'lt2',
    'TEXT_1': 'dk1', 'DARK_1': 'dk1',
    'TEXT_2': 'dk2', 'DARK_2': 'dk2',
    'ACCENT_1': 'accent1', 'ACCENT_2': 'accent2',
    'ACCENT_3': 'accent3', 'ACCENT_4': 'accent4',
    'ACCENT_5': 'accent5', 'ACCENT_6': 'accent6',
    'HYPERLINK': 'hlink', 'FOLLOWED_HYPERLINK': 'folHlink',
}


def emu_to_cm(emu: int) -> float:
    return round(emu / EMU_PER_CM, 2)


def emu_to_px(emu: int, dpi: int = 96) -> int:
    return int(emu / EMU_PER_INCH * dpi)


def emu_to_pt(emu: int) -> float:
    return round(emu / EMU_PER_PT, 1)


def emu_to_percent(emu: int, total_emu: int) -> float:
    if total_emu == 0:
        return 0
    return round((emu / total_emu) * 100, 2)


def safe_rgb(color_obj, theme_colors: Optional[Dict[str, str]] = None) -> Optional[str]:
    """安全获取颜色 RGB，支持通过 theme_colors 解析 schemeClr"""
    try:
        return f"#{str(color_obj.rgb).upper()}"
    except Exception:
        pass
    if theme_colors:
        try:
            from pptx.enum.dml import MSO_COLOR_TYPE
            if color_obj.type == MSO_COLOR_TYPE.SCHEME:
                key = _TC_KEY_MAP.get(color_obj.theme_color.name)
                if key and key in theme_colors:
                    return theme_colors[key]
        except Exception:
            pass
    return None


# ─── Theme XML 解析（zipfile 直读，无需 python-pptx）──────────────────────────


def parse_color(color_elem) -> Optional[str]:
    if color_elem is None:
        return None
    srgb = color_elem.find(f'.//{{{NS}}}srgbClr')
    if srgb is not None:
        return f"#{srgb.get('val', '000000').upper()}"
    sys_clr = color_elem.find(f'.//{{{NS}}}sysClr')
    if sys_clr is not None:
        last = sys_clr.get('lastClr')
        if last:
            return f"#{last.upper()}"
        return {'windowText': '#000000', 'window': '#FFFFFF'}.get(sys_clr.get('val', ''), '#000000')
    return None


def extract_theme_colors(theme_xml: str) -> Dict[str, str]:
    import xml.etree.ElementTree as ET
    colors = {}
    try:
        root = ET.fromstring(theme_xml)
        ns = {'a': NS}
        scheme = root.find('.//a:clrScheme', ns)
        if scheme is not None:
            for name in ['dk1', 'lt1', 'dk2', 'lt2',
                         'accent1', 'accent2', 'accent3', 'accent4', 'accent5', 'accent6',
                         'hlink', 'folHlink']:
                elem = scheme.find(f'a:{name}', ns)
                if elem is not None:
                    c = parse_color(elem)
                    if c:
                        colors[name] = c
    except Exception as e:
        print(f"解析颜色方案失败: {e}", file=sys.stderr)
    return colors


def extract_theme_fonts(theme_xml: str) -> Dict[str, Any]:
    import xml.etree.ElementTree as ET
    fonts = {}
    try:
        root = ET.fromstring(theme_xml)
        ns = {'a': NS}
        scheme = root.find('.//a:fontScheme', ns)
        if scheme is not None:
            for role, key in [('majorFont', 'major'), ('minorFont', 'minor')]:
                elem = scheme.find(f'.//a:{role}', ns)
                if elem is not None:
                    fonts[key] = {
                        ft: elem.find(f'a:{ft}', ns).get('typeface', '')
                        for ft in ['latin', 'ea', 'cs']
                        if elem.find(f'a:{ft}', ns) is not None
                    }
    except Exception as e:
        print(f"解析字体方案失败: {e}", file=sys.stderr)
    return fonts


def extract_slide_size(pres_xml: str) -> Dict[str, Any]:
    import xml.etree.ElementTree as ET
    try:
        root = ET.fromstring(pres_xml)
        ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
        sz = root.find('.//p:sldSz', ns)
        if sz is not None:
            cx = int(sz.get('cx', '9144000'))
            cy = int(sz.get('cy', '6858000'))
            return {
                'width_emu': cx, 'height_emu': cy,
                'width_cm': emu_to_cm(cx), 'height_cm': emu_to_cm(cy),
                'width_px': emu_to_px(cx), 'height_px': emu_to_px(cy),
                'aspect_ratio': f"{cx // 10000}:{cy // 10000}"
            }
    except Exception as e:
        print(f"解析幻灯片尺寸失败: {e}", file=sys.stderr)
    return {}


# ─── 形状级属性提取（python-pptx + 直读 XML）─────────────────────────────────


def extract_border(shape) -> Dict[str, Any]:
    """提取形状边框：粗细(pt)、颜色、线型"""
    out = {}
    try:
        line = shape.line
        if line.width is not None:
            out['width_pt'] = emu_to_pt(line.width)
        color = safe_rgb(line.color)
        if color:
            out['color'] = color
        if line.dash_style is not None:
            out['dash'] = str(line.dash_style)
    except Exception:
        pass
    return out


def extract_shadow(shape) -> Dict[str, Any]:
    """从 XML 提取阴影参数：模糊半径、偏移距离、方向、颜色、透明度"""
    out = {}
    try:
        outer = shape._element.find(f'.//{{{NS}}}outerShdw')
        if outer is None:
            return out
        for attr, key, scale in [
            ('blurRad', 'blur_pt', EMU_PER_PT),
            ('dist',    'distance_pt', EMU_PER_PT),
        ]:
            val = outer.get(attr)
            if val:
                out[key] = round(int(val) / scale, 1)
        direction = outer.get('dir')
        if direction:
            out['direction_deg'] = round(int(direction) / 60000, 1)
        clr = outer.find(f'{{{NS}}}srgbClr')
        if clr is not None:
            out['color'] = f"#{clr.get('val', '000000').upper()}"
            alpha = clr.find(f'{{{NS}}}alpha')
            if alpha is not None:
                out['opacity_pct'] = round(int(alpha.get('val', '100000')) / 1000, 1)
    except Exception:
        pass
    return out


def extract_corner(shape) -> Optional[str]:
    """提取圆角类型及比例（roundRect 返回百分比，其他几何体返回名称）"""
    try:
        geom = shape._element.find(f'.//{{{NS}}}prstGeom')
        if geom is None:
            return None
        prst = geom.get('prst', '')
        if prst == 'roundRect':
            av_lst = geom.find(f'{{{NS}}}avLst')
            if av_lst is not None:
                gd = av_lst.find(f'{{{NS}}}gd')
                if gd is not None:
                    m = re.search(r'val (\d+)', gd.get('fmla', ''))
                    if m:
                        return f"{round(int(m.group(1)) / 100000 * 100, 1)}%"
            return 'roundRect'
        return prst if prst not in ('rect', '') else None
    except Exception:
        return None


def extract_gradient(shape) -> Optional[Dict[str, Any]]:
    """提取渐变填充：stops（颜色+位置）和角度"""
    try:
        grad = shape._element.find(f'.//{{{NS}}}gradFill')
        if grad is None:
            return None
        result: Dict[str, Any] = {'stops': []}
        lin = grad.find(f'.//{{{NS}}}lin')
        if lin is not None and lin.get('ang'):
            result['angle_deg'] = round(int(lin.get('ang')) / 60000, 1)
        gs_lst = grad.find(f'.//{{{NS}}}gsLst')
        if gs_lst is not None:
            for gs in gs_lst.findall(f'{{{NS}}}gs'):
                stop: Dict[str, Any] = {}
                pos = gs.get('pos')
                if pos:
                    stop['position_pct'] = round(int(pos) / 1000, 1)
                srgb = gs.find(f'{{{NS}}}srgbClr')
                if srgb is not None:
                    stop['color'] = f"#{srgb.get('val', '').upper()}"
                    alpha = srgb.find(f'{{{NS}}}alpha')
                    if alpha is not None:
                        stop['opacity_pct'] = round(int(alpha.get('val', '100000')) / 1000, 1)
                if stop:
                    result['stops'].append(stop)
        return result if result['stops'] else None
    except Exception:
        return None


def _collect_gradient_stop_colors(elem) -> List[str]:
    """从任意 XML 元素中递归收集 gradFill 的 srgbClr 颜色列表"""
    colors = []
    for grad in elem.iter(f'{{{NS}}}gradFill'):
        gs_lst = grad.find(f'{{{NS}}}gsLst')
        if gs_lst is None:
            continue
        for gs in gs_lst.findall(f'{{{NS}}}gs'):
            srgb = gs.find(f'{{{NS}}}srgbClr')
            if srgb is not None:
                val = srgb.get('val', '')
                if val:
                    colors.append(f"#{val.upper()}")
    return colors


# ─── 幻灯片级别扫描（python-pptx）────────────────────────────────────────────

# placeholder idx -> 语义标签
_PH_LABEL = {
    0: 'title', 1: 'body', 2: 'subtitle',
    3: 'center_title', 10: 'date', 11: 'footer', 12: 'slide_number',
}


def _shape_role(shape) -> str:
    if shape.is_placeholder:
        idx = shape.placeholder_format.idx
        return _PH_LABEL.get(idx, f'placeholder_{idx}')
    return 'free_shape'


def resolve_attr(getters: list) -> Any:
    """依次执行 getters，返回第一个非 None 值；任何异常自动跳过"""
    for getter in getters:
        try:
            val = getter()
            if val is not None:
                return val
        except Exception:
            pass
    return None


def _find_matching_ph(shape, container) -> Optional[Any]:
    """在 layout 或 master 中找与 shape 同 idx 的占位符"""
    if not shape.is_placeholder:
        return None
    idx = shape.placeholder_format.idx
    for ph in container.placeholders:
        try:
            if ph.placeholder_format.idx == idx:
                return ph
        except Exception:
            pass
    return None


def _para_defrpr_sz(para) -> Optional[float]:
    """从段落 pPr/defRPr 读字号（centi-pt → pt）"""
    pPr = para._p.find(f'{{{NS}}}pPr')
    if pPr is None:
        return None
    defrpr = pPr.find(f'{{{NS}}}defRPr')
    if defrpr is None:
        return None
    val = defrpr.get('sz')
    return int(val) / 100 if val else None


def _ph_defrpr_size_pt(ph) -> Optional[float]:
    """从占位符第一个段落的 defRPr 读字号"""
    for para in ph.text_frame.paragraphs:
        sz = _para_defrpr_sz(para)
        if sz:
            return sz
    return None


def _ph_defrpr_font(ph) -> Optional[str]:
    """从占位符第一个段落的 defRPr latin 读字体名"""
    for para in ph.text_frame.paragraphs:
        pPr = para._p.find(f'{{{NS}}}pPr')
        if pPr is None:
            continue
        defrpr = pPr.find(f'{{{NS}}}defRPr')
        if defrpr is None:
            continue
        latin = defrpr.find(f'{{{NS}}}latin')
        if latin is not None:
            tf = latin.get('typeface', '')
            if tf:
                return tf
    return None


def _resolve_theme_font(typeface: str, theme_fonts: Dict[str, Any]) -> str:
    """将 +mj-lt/+mn-lt/+mj-ea/+mn-ea 等主题字体引用解析为实际字体名"""
    mapping = {
        '+mj-lt': theme_fonts.get('major', {}).get('latin', ''),
        '+mn-lt': theme_fonts.get('minor', {}).get('latin', ''),
        '+mj-ea': theme_fonts.get('major', {}).get('ea', ''),
        '+mn-ea': theme_fonts.get('minor', {}).get('ea', ''),
    }
    return mapping.get(typeface, typeface) or typeface


def _extract_tx_styles_defaults(prs) -> Dict[str, Any]:
    """从 master txStyles 提取标题/正文的默认字号和字体（作为继承链最终兜底）"""
    result: Dict[str, Any] = {}
    try:
        tx = prs.slide_master._element.find(f'{{{NS_P_MAIN}}}txStyles')
        if tx is None:
            return result
        for style_tag, role in [('titleStyle', 'title'), ('bodyStyle', 'body')]:
            elem = tx.find(f'{{{NS_P_MAIN}}}{style_tag}')
            if elem is None:
                continue
            lvl1 = elem.find(f'.//{{{NS}}}lvl1pPr')
            if lvl1 is None:
                continue
            defrpr = lvl1.find(f'{{{NS}}}defRPr')
            if defrpr is None:
                continue
            sz = defrpr.get('sz')
            if sz:
                result[f'{role}_size_pt'] = int(sz) / 100
            latin = defrpr.find(f'{{{NS}}}latin')
            if latin is not None:
                tf = latin.get('typeface', '')
                if tf:
                    result[f'{role}_font'] = tf
    except Exception:
        pass
    return result


def _resolve_bg_color(slide, theme_colors: Optional[Dict[str, str]] = None) -> Optional[str]:
    """按 slide → layout → master 优先级解析背景纯色，找到第一个有效定义即返回"""
    for obj in (slide, slide.slide_layout, slide.slide_layout.slide_master):
        try:
            fill = obj.background.fill
            if fill.type is not None:
                c = safe_rgb(fill.fore_color, theme_colors)
                if c:
                    return c
        except Exception:
            pass
    return None


def extract_actual_colors(prs, theme_colors: Optional[Dict[str, str]] = None) -> Dict[str, Any]:
    """扫描所有幻灯片，统计实际使用颜色（含非主题色）及用途。
    按面积加权排序：大形状填充色权重高，文字色单独计数，不与填充混排。
    """
    stats: Dict[str, Dict] = defaultdict(lambda: {
        'count': 0,
        'fill_count': 0,   # 作为填充/背景出现的次数
        'text_count': 0,   # 作为文字色出现的次数
        'area_weight': 0.0,  # 填充面积权重之和（背景=1.0，形状=shape面积/幻灯片面积）
        'usages': set(),
    })

    slide_area = prs.slide_width * prs.slide_height

    def record(color: Optional[str], usage: str, area_pct: float = 0.0):
        if not color:
            return
        stats[color]['count'] += 1
        stats[color]['usages'].add(usage)
        if '_text' in usage:
            stats[color]['text_count'] += 1
        else:
            stats[color]['fill_count'] += 1
            stats[color]['area_weight'] += area_pct

    for slide in prs.slides:
        # 幻灯片背景色（slide → layout → master 优先级），面积权重为 1.0
        c = _resolve_bg_color(slide, theme_colors)
        if c:
            record(c, 'background_fill', 1.0)

        for shape in slide.shapes:
            role = _shape_role(shape)
            # 形状面积占幻灯片比例（限制在 [0, 1]，防止超出边界的形状权重异常）
            try:
                shape_area_pct = min(1.0, (shape.width * shape.height) / slide_area) if slide_area else 0.0
            except Exception:
                shape_area_pct = 0.0

            # 填充色（纯色）
            try:
                c = safe_rgb(shape.fill.fore_color, theme_colors)
                record(c, f'{role}_fill', shape_area_pct)
            except Exception:
                pass
            # 填充色（渐变 stops）：渐变各 stop 共享形状面积权重
            # pylint: disable=protected-access
            grad_colors = _collect_gradient_stop_colors(shape._element.find(f'.//{{{NS}}}spPr') or shape._element)
            # pylint: enable=protected-access
            stop_pct = shape_area_pct / len(grad_colors) if grad_colors else 0.0
            for c in grad_colors:
                record(c, f'{role}_fill_gradient', stop_pct)
            # 边框色：视觉面积极小，给固定微权重
            try:
                c = safe_rgb(shape.line.color, theme_colors)
                record(c, f'{role}_border', 0.005)
            except Exception:
                pass
            # 文字色（纯色）：text_count 计次，area_weight 不增加
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        try:
                            c = safe_rgb(run.font.color, theme_colors)
                            record(c, f'{role}_text', 0.0)
                        except Exception:
                            pass
                # 文字渐变色（gradFill 在 rPr 内）
                for rPr in shape._element.iter(f'{{{NS}}}rPr'):
                    for c in _collect_gradient_stop_colors(rPr):
                        record(c, f'{role}_text_gradient', 0.0)

    # 排序：优先按 area_weight 降序（视觉面积主导），面积相同时按 fill_count 降序
    sorted_items = sorted(
        stats.items(),
        key=lambda x: (-x[1]['area_weight'], -x[1]['fill_count'], -x[1]['count'])
    )
    return {
        color: {
            'count': d['count'],
            'fill_count': d['fill_count'],
            'text_count': d['text_count'],
            'area_weight': round(d['area_weight'], 4),
            'usages': sorted(d['usages']),
        }
        for color, d in sorted_items[:24]
    }


def extract_bg_text_mapping(prs, theme_colors: Optional[Dict[str, str]] = None) -> Dict[str, Any]:
    """
    统计每种背景色上实际使用的文字色，输出背景→文字配色映射。
    背景色按 slide → layout → master 优先级解析，覆盖继承背景场景。
    """
    bg_to_texts: Dict[str, Counter] = defaultdict(Counter)
    bg_slide_count: Dict[str, int] = defaultdict(int)

    for slide in prs.slides:
        bg_color = _resolve_bg_color(slide, theme_colors)

        if not bg_color:
            continue

        bg_slide_count[bg_color] += 1

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    try:
                        tc = safe_rgb(run.font.color, theme_colors)
                        if tc:
                            bg_to_texts[bg_color][tc] += 1
                    except Exception:
                        pass

    result = {}
    for bg, text_counter in bg_to_texts.items():
        if text_counter:
            result[bg] = {
                'text_colors': sorted(text_counter.keys(), key=lambda c: -text_counter[c]),
                'slide_count': bg_slide_count[bg],
            }
    return result


def extract_font_sizes(prs, theme_fonts: Optional[Dict[str, Any]] = None) -> Dict[str, Any]:
    """
    按占位符层级统计实际字号和字体名。
    字号/字体继承链：run → para defRPr → layout ph defRPr → master ph defRPr → txStyles
    """
    size_buckets: Dict[str, List[float]] = defaultdict(list)
    font_buckets: Dict[str, List[str]]  = defaultdict(list)

    tx_defaults = _extract_tx_styles_defaults(prs)
    tf = theme_fonts or {}

    for slide in prs.slides:
        layout = slide.slide_layout
        master = layout.slide_master

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            role = _shape_role(shape)
            layout_ph = _find_matching_ph(shape, layout)
            master_ph  = _find_matching_ph(shape, master)
            tx_sz_key  = 'title_size_pt' if 'title' in role else 'body_size_pt'
            tx_fn_key  = 'title_font'    if 'title' in role else 'body_font'

            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    # ── 字号 ──
                    para_sz = _para_defrpr_sz(para)
                    size_emu = resolve_attr([
                        lambda: run.font.size,
                        lambda: para_sz * EMU_PER_PT if para_sz else None,
                    ])
                    if size_emu:
                        size_buckets[role].append(emu_to_pt(size_emu))
                    else:
                        default_pt = resolve_attr([
                            lambda: _ph_defrpr_size_pt(layout_ph) if layout_ph else None,
                            lambda: _ph_defrpr_size_pt(master_ph) if master_ph else None,
                            lambda: tx_defaults.get(tx_sz_key),
                        ])
                        if default_pt:
                            size_buckets[role].append(default_pt)

                    # ── 字体名 ──
                    raw = resolve_attr([
                        lambda: run.font.name or None,
                        lambda: _ph_defrpr_font(layout_ph) if layout_ph else None,
                        lambda: _ph_defrpr_font(master_ph) if master_ph else None,
                        lambda: tx_defaults.get(tx_fn_key),
                    ])
                    if raw:
                        resolved = _resolve_theme_font(raw, tf)
                        if resolved and not resolved.startswith('+'):
                            font_buckets[role].append(resolved)

    result = {}
    for role in set(size_buckets) | set(font_buckets):
        sizes = size_buckets.get(role, [])
        fonts = font_buckets.get(role, [])
        entry: Dict[str, Any] = {}
        if sizes:
            top = Counter(sizes).most_common(3)
            entry['dominant_pt']      = top[0][0]
            entry['common_sizes_pt']  = [s for s, _ in top]
            entry['sample_count']     = len(sizes)
        if fonts:
            entry['dominant_font'] = Counter(fonts).most_common(1)[0][0]
        if entry:
            result[role] = entry
    return result


def extract_para_alignment(prs) -> Dict[str, Any]:
    """
    按占位符层级统计段落对齐方式。
    对齐继承链：para.alignment → layout ph 首段 → master ph 首段
    """
    align_buckets: Dict[str, Counter] = defaultdict(Counter)

    def _ph_align(ph) -> Optional[str]:
        try:
            a = ph.text_frame.paragraphs[0].alignment
            return str(a) if a is not None else None
        except Exception:
            return None

    for slide in prs.slides:
        layout = slide.slide_layout
        master = layout.slide_master

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            role = _shape_role(shape)
            layout_ph = _find_matching_ph(shape, layout)
            master_ph  = _find_matching_ph(shape, master)

            for para in shape.text_frame.paragraphs:
                align = resolve_attr([
                    lambda: str(para.alignment) if para.alignment is not None else None,
                    lambda: _ph_align(layout_ph) if layout_ph else None,
                    lambda: _ph_align(master_ph) if master_ph else None,
                ])
                if align:
                    align_buckets[role][align] += 1

    result = {}
    for role, counts in align_buckets.items():
        if counts:
            result[role] = {
                'dominant': counts.most_common(1)[0][0],
                'counts': dict(counts),
            }
    return result


def extract_master_elements(prs, theme_colors: Optional[Dict[str, str]] = None) -> Dict[str, Any]:
    """提取母版（Slide Master）固定装饰元素：背景、形状列表"""
    master = prs.slide_master
    w, h = prs.slide_width, prs.slide_height
    info: Dict[str, Any] = {'background': {}, 'fixed_shapes': [], 'placeholders': []}

    # 母版背景
    try:
        c = safe_rgb(master.background.fill.fore_color, theme_colors)
        if c:
            info['background'] = {'type': 'solid', 'color': c}
    except Exception:
        pass

    for shape in master.shapes:
        entry: Dict[str, Any] = {
            'shape_type': shape.shape_type.name if hasattr(shape.shape_type, 'name') else str(shape.shape_type),
            'left_pct': emu_to_percent(shape.left, w),
            'top_pct':  emu_to_percent(shape.top, h),
            'width_pct': emu_to_percent(shape.width, w),
            'height_pct': emu_to_percent(shape.height, h),
            'left_px': emu_to_px(shape.left),
            'top_px':  emu_to_px(shape.top),
            'width_px': emu_to_px(shape.width),
            'height_px': emu_to_px(shape.height),
        }
        try:
            c = safe_rgb(shape.fill.fore_color, theme_colors)
            if c:
                entry['fill_color'] = c
        except Exception:
            pass
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                entry['text_preview'] = text[:60]
        border = extract_border(shape)
        if border:
            entry['border'] = border

        if shape.is_placeholder:
            entry['placeholder_idx'] = shape.placeholder_format.idx
            info['placeholders'].append(entry)
        else:
            info['fixed_shapes'].append(entry)

    return info


def extract_component_styles(prs) -> Dict[str, Any]:
    """扫描非占位符形状，统计边框/阴影/圆角/渐变的典型值"""
    borders, shadows, corners, gradients = [], [], [], []

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.is_placeholder:
                continue
            b = extract_border(shape)
            if b:
                borders.append(b)
            s = extract_shadow(shape)
            if s:
                shadows.append(s)
            c = extract_corner(shape)
            if c:
                corners.append(c)
            g = extract_gradient(shape)
            if g:
                gradients.append(g)

    def dominant(lst, key):
        vals = [x[key] for x in lst if key in x]
        return Counter(vals).most_common(1)[0][0] if vals else None

    result: Dict[str, Any] = {}
    if borders:
        result['border'] = {
            'typical_width_pt': dominant(borders, 'width_pt'),
            'typical_color': dominant(borders, 'color'),
            'count': len(borders),
        }
    if shadows:
        result['shadow'] = {
            'typical_blur_pt': dominant(shadows, 'blur_pt'),
            'typical_distance_pt': dominant(shadows, 'distance_pt'),
            'typical_color': dominant(shadows, 'color'),
            'count': len(shadows),
        }
    if corners:
        top_corner = Counter(corners).most_common(1)[0][0]
        result['corner'] = {
            'typical': top_corner,
            'variants': list(Counter(corners).keys()),
            'count': len(corners),
        }
    if gradients:
        result['gradient'] = {
            'count': len(gradients),
            'examples': gradients[:3],
        }
    return result


# ─── 版式提取 ────────────────────────────────────────────────────────────────


def extract_layouts(prs, theme_fonts: Optional[Dict[str, Any]] = None) -> List[Dict[str, Any]]:
    """
    提取所有版式的占位符信息，含字号/字体/对齐。
    占位符位置若 layout 层未定义则向 master 回退；字号/字体/对齐同样走继承链。
    """
    layouts = []
    try:
        w, h = prs.slide_width, prs.slide_height
        master = prs.slide_master
        tx_defaults = _extract_tx_styles_defaults(prs)
        tf = theme_fonts or {}

        for i, layout in enumerate(prs.slide_layouts):
            info: Dict[str, Any] = {'index': i, 'name': layout.name, 'placeholders': []}
            for shape in layout.shapes:
                if not shape.is_placeholder:
                    continue
                master_ph = _find_matching_ph(shape, master)
                idx = shape.placeholder_format.idx
                role = _PH_LABEL.get(idx, f'placeholder_{idx}')
                tx_sz_key = 'title_size_pt' if 'title' in role else 'body_size_pt'
                tx_fn_key = 'title_font'    if 'title' in role else 'body_font'

                # 位置：layout 优先，layout 值为 0 且 master 有定义时回退
                def _pos(attr, fallback_obj):
                    v = getattr(shape, attr, None)
                    if v:
                        return v
                    return getattr(fallback_obj, attr, None) if fallback_obj else None

                left   = _pos('left',   master_ph) or 0
                top    = _pos('top',    master_ph) or 0
                width  = _pos('width',  master_ph) or 0
                height = _pos('height', master_ph) or 0

                ph: Dict[str, Any] = {
                    'type': str(shape.placeholder_format.type),
                    'idx': idx,
                    'left_cm':      emu_to_cm(left),
                    'top_cm':       emu_to_cm(top),
                    'width_cm':     emu_to_cm(width),
                    'height_cm':    emu_to_cm(height),
                    'left_percent':   emu_to_percent(left,   w),
                    'top_percent':    emu_to_percent(top,    h),
                    'width_percent':  emu_to_percent(width,  w),
                    'height_percent': emu_to_percent(height, h),
                }

                # 字号
                sz = resolve_attr([
                    lambda: _ph_defrpr_size_pt(shape) if shape.has_text_frame else None,
                    lambda: _ph_defrpr_size_pt(master_ph) if master_ph else None,
                    lambda: tx_defaults.get(tx_sz_key),
                ])
                if sz:
                    ph['font_size_pt'] = sz

                # 字体名
                raw = resolve_attr([
                    lambda: _ph_defrpr_font(shape) if shape.has_text_frame else None,
                    lambda: _ph_defrpr_font(master_ph) if master_ph else None,
                    lambda: tx_defaults.get(tx_fn_key),
                ])
                if raw:
                    resolved = _resolve_theme_font(raw, tf)
                    if resolved and not resolved.startswith('+'):
                        ph['font_name'] = resolved

                # 对齐
                align = resolve_attr([
                    lambda: (str(shape.text_frame.paragraphs[0].alignment)
                             if shape.has_text_frame and shape.text_frame.paragraphs
                                and shape.text_frame.paragraphs[0].alignment is not None
                             else None),
                    lambda: (str(master_ph.text_frame.paragraphs[0].alignment)
                             if master_ph and master_ph.has_text_frame
                                and master_ph.text_frame.paragraphs
                                and master_ph.text_frame.paragraphs[0].alignment is not None
                             else None),
                ])
                if align:
                    ph['alignment'] = align

                info['placeholders'].append(ph)
            layouts.append(info)
    except Exception as e:
        print(f"提取版式信息失败: {e}", file=sys.stderr)
    return layouts


def _shape_text(shape) -> str:
    try:
        if shape.has_text_frame:
            return shape.text_frame.text.strip()
    except Exception:
        pass
    return ''


def _shape_box_pct(shape, slide_w: int, slide_h: int) -> Dict[str, float]:
    return {
        'left': emu_to_percent(getattr(shape, 'left', 0), slide_w),
        'top': emu_to_percent(getattr(shape, 'top', 0), slide_h),
        'width': emu_to_percent(getattr(shape, 'width', 0), slide_w),
        'height': emu_to_percent(getattr(shape, 'height', 0), slide_h),
    }


def _bucket(value: float, step: int = 10) -> int:
    return int(round(value / step) * step)


def _is_visual_shape(shape) -> bool:
    try:
        shape_type = getattr(shape.shape_type, 'name', str(shape.shape_type))
    except Exception:
        shape_type = ''
    if 'PICTURE' in shape_type or 'CHART' in shape_type or 'TABLE' in shape_type:
        return True
    if 'GROUP' in shape_type or 'FREEFORM' in shape_type or 'LINE' in shape_type:
        return False
    try:
        if shape.fill and shape.fill.type is not None and not _shape_text(shape):
            return True
    except Exception:
        pass
    return False


def _collect_slide_features(slide, slide_w: int, slide_h: int) -> Dict[str, Any]:
    texts = []
    text_boxes = []
    visuals = []
    thin_lines = []

    for shape in slide.shapes:
        try:
            shape_type = getattr(shape.shape_type, 'name', str(shape.shape_type))
        except Exception:
            shape_type = ''
        text = _shape_text(shape)
        box = _shape_box_pct(shape, slide_w, slide_h)
        area = box['width'] * box['height']
        if text:
            font_size = None
            try:
                sizes = []
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size:
                            sizes.append(run.font.size.pt)
                if sizes:
                    font_size = max(sizes)
            except Exception:
                pass
            text_boxes.append({'text': text, 'box': box, 'area': area, 'font_size': font_size})
            texts.append(text)

        is_auto_shape_visual = 'AUTO_SHAPE' in shape_type and not text and area >= 80
        if _is_visual_shape(shape) or is_auto_shape_visual:
            visuals.append({'box': box, 'area': area, 'shape_type': shape_type})

        is_horizontal_line = box['width'] >= 18 and box['height'] <= 2.0
        is_vertical_line = box['height'] >= 18 and box['width'] <= 2.0
        if 'LINE' in shape_type and (is_horizontal_line or is_vertical_line):
            thin_lines.append({'box': box, 'area': area})

    all_text = '\n'.join(texts)
    title_like = []
    for t in text_boxes:
        is_large_font = (t.get('font_size') or 0) >= 28
        is_top_positioned = t['box']['top'] <= 12 and t['box']['height'] <= 10
        if is_large_font or is_top_positioned:
            title_like.append(t)
    body_boxes = []
    for t in text_boxes:
        is_body = t not in title_like and len(t['text'].strip()) >= 2 and t['area'] >= 8
        if is_body:
            body_boxes.append(t)
    large_visuals = [v for v in visuals if v['area'] >= 30]
    return {
        'text': all_text,
        'text_boxes': text_boxes,
        'title_like': title_like,
        'body_boxes': body_boxes,
        'visuals': visuals,
        'large_visuals': large_visuals,
        'thin_lines': thin_lines,
    }


def _content_subtype(features: Dict[str, Any]) -> str:
    body_count = len(features['body_boxes'])
    visual_count = len(features['large_visuals'])
    thin_count = len(features['thin_lines'])

    if thin_count and body_count >= 3:
        return 'content-timeline'
    if visual_count >= 3 or body_count >= 4:
        return 'content-multi-panel'
    if visual_count >= 1 and body_count >= 1:
        return 'content-text-image'
    if body_count >= 2:
        return 'content-two-column-text'
    return 'content-general'


def classify_slide_roles(prs) -> List[Dict[str, Any]]:
    """基于 PPTX 源文件结构识别封面、目录、章节、内容和结束页。"""
    slide_w, slide_h = prs.slide_width, prs.slide_height
    roles: List[Dict[str, Any]] = []
    total = len(prs.slides)

    for idx, slide in enumerate(prs.slides, start=1):
        features = _collect_slide_features(slide, slide_w, slide_h)
        text = features['text']
        compact = re.sub(r'\s+', '', text).upper()
        text_count = len(features['text_boxes'])
        body_count = len(features['body_boxes'])
        role = 'content'
        subtype = _content_subtype(features)
        confidence = 0.62
        evidence = []

        numeric_markers = [
            t for t in features['text_boxes']
            if re.fullmatch(r'0?[1-9]', re.sub(r'\s+', '', t['text']))
        ]
        part_markers = []
        for t in features['text_boxes']:
            is_part_marker = (
                re.search(r'PART\s*0?[1-9]', t['text'], re.IGNORECASE)
                or re.search(r'第[一二三四五六七八九十0-9]+章节?', t['text'])
                or re.search(r'0?[1-9]\s*[、.．-]?\s*章节', t['text'])
            )
            if is_part_marker:
                part_markers.append(t)
        has_toc_keyword = bool(re.search(r'(目录|CONTENTS|AGENDA)', compact))
        has_cover_signal = bool(re.search(r'(汇报人|模板|工作总结|年终总结|述职报告|YOURLOGO|BUSINESS|REPORT)', compact))
        has_cover_text_density = text_count <= 8 or has_cover_signal

        if re.search(r'(THANK\s*YOU|谢谢|致谢)', compact):
            role, subtype, confidence = 'closing', 'closing-thanks', 0.92
            evidence.append('包含结束页关键词')
        elif has_toc_keyword:
            role, subtype, confidence = 'toc', 'toc-numbered-list', 0.9
            evidence.append('包含目录关键词')
        elif idx == 1 and body_count <= 2 and has_cover_text_density:
            role, subtype, confidence = 'cover', 'cover-title-subtitle', 0.8
            evidence.append('首页且具有封面标题/汇报信息特征')
        elif (numeric_markers or part_markers) and idx != 1:
            role, subtype, confidence = 'section', 'section-number-title', 0.82
            evidence.append('包含章节编号或 PART 标记')
        elif idx == 1 and text_count <= 5 and body_count <= 1:
            role, subtype, confidence = 'cover', 'cover-title-subtitle', 0.78
            evidence.append('首页且文本密度低')
        else:
            evidence.append(f'正文块 {body_count} 个，图片区 {len(features["large_visuals"])} 个')

        roles.append({
            'page': idx,
            'role': role,
            'subtype': subtype,
            'confidence': round(confidence, 2),
            'evidence': evidence,
        })

    return roles


def _region_name(box: Dict[str, float]) -> str:
    cx = box['left'] + box['width'] / 2
    cy = box['top'] + box['height'] / 2
    x = 'left' if cx < 40 else 'right' if cx > 60 else 'center'
    y = 'top' if cy < 35 else 'bottom' if cy > 65 else 'middle'
    return f'{x}-{y}'


def _content_layout_signature(features: Dict[str, Any], subtype: str) -> str:
    body_regions = sorted(_region_name(t['box']) for t in features['body_boxes'])
    visual_regions = sorted(_region_name(v['box']) for v in features['large_visuals'])
    return json.dumps({
        'subtype': subtype,
        'body_regions': body_regions,
        'visual_regions': visual_regions,
        'body_count': len(features['body_boxes']),
        'visual_count': len(features['large_visuals']),
        'thin_lines': len(features['thin_lines']),
    }, ensure_ascii=False, sort_keys=True)


def _describe_content_style(subtype: str, features: Dict[str, Any]) -> str:
    bodies = features['body_boxes']
    visuals = features['large_visuals']
    if subtype == 'content-timeline':
        return '时间轴/流程页：用横向或纵向细线串联多个节点，节点旁放短文本说明。'
    if subtype == 'content-multi-panel':
        return f'多模块内容页：约 {len(bodies)} 个正文/说明模块，配合 {len(visuals)} 个图片区或装饰图片区。'
    if subtype == 'content-text-image' and len(bodies) >= 2 and visuals:
        v_region = _region_name(visuals[0]['box']).replace('-', ' ')
        return f'双段正文 + 配图：正文分成 {len(bodies)} 个大文本块，图片区位于 {v_region} 区域。'
    if subtype == 'content-text-image':
        return '图文混排内容页：正文说明与一个主要图片区形成左右或上下对照。'
    if subtype == 'content-two-column-text':
        return f'多段正文页：{len(bodies)} 个大文本块分区排布，适合并列说明或分层介绍。'
    return '通用内容页：标题下组织正文、说明或少量图形元素。'


def _usage_rule(subtype: str) -> str:
    rules = {
        'content-text-image': '适合包含说明文字、案例/照片/图示的内容；按图文比例选择左右、上下或角落配图。',
        'content-timeline': '适合发展历程、流程步骤、阶段规划等有顺序关系的内容。',
        'content-two-column-text': '适合纯文字信息较多、需要分组对照的内容。',
        'content-multi-panel': '适合多对象展示、多维度对比、校园风采/案例集等图片或模块较多的内容。',
        'content-general': '适合信息量较少或无法匹配其他专用样式的普通内容页。',
    }
    return rules.get(subtype, rules['content-general'])


def _semantic_guess(subtype: str, features: Dict[str, Any]) -> Dict[str, str]:
    bodies = features['body_boxes']
    visuals = features['large_visuals']
    thin_lines = features['thin_lines']
    body_count = len(bodies)
    visual_count = len(visuals)

    if subtype == 'content-timeline' or (thin_lines and body_count >= 3):
        return {
            'semantic_type_guess': 'process',
            'information_relation': '多个节点按时间、步骤或阶段顺序推进。',
            'selection_rule': '适用于流程步骤、阶段规划、时间线和路径推进类内容。',
        }
    if subtype == 'content-text-image' or (visual_count >= 1 and body_count >= 1):
        return {
            'semantic_type_guess': 'image_text',
            'information_relation': '图片提供场景、证据或视觉说明，文字解释核心观点。',
            'selection_rule': '适用于单图说明、案例展示、场景介绍和图文解读。',
        }
    if subtype == 'content-multi-panel' or body_count >= 4 or visual_count >= 3:
        return {
            'semantic_type_guess': 'parallel',
            'information_relation': '多个模块并列呈现，同属一个主题但彼此独立。',
            'selection_rule': '适用于多对象展示、多维度对比、功能模块和案例集合。',
        }
    if subtype == 'content-two-column-text':
        left_count = sum(1 for b in bodies if _region_name(b['box']).startswith('left'))
        right_count = sum(1 for b in bodies if _region_name(b['box']).startswith('right'))
        semantic = 'comparison' if left_count and right_count else 'parallel'
        return {
            'semantic_type_guess': semantic,
            'information_relation': '两组或多组文本分区表达，可用于对比、并列说明或分层展开。',
            'selection_rule': '适用于纯文字信息较多、需要分组对照或并列展开的内容。',
        }
    return {
        'semantic_type_guess': 'unknown',
        'information_relation': '信息关系不明显，按通用内容页处理。',
        'selection_rule': '适用于信息量较少或无法匹配专用版式的普通内容。',
    }


def extract_content_layout_styles(prs, slide_roles: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    """抽取内容页版式样式并按结构指纹去重。"""
    slide_w, slide_h = prs.slide_width, prs.slide_height
    role_by_page = {r['page']: r for r in slide_roles}
    styles_by_sig: Dict[str, Dict[str, Any]] = {}

    for idx, slide in enumerate(prs.slides, start=1):
        role = role_by_page.get(idx, {})
        if role.get('role') != 'content':
            continue
        features = _collect_slide_features(slide, slide_w, slide_h)
        subtype = role.get('subtype') or _content_subtype(features)
        signature = _content_layout_signature(features, subtype)
        if signature not in styles_by_sig:
            semantic = _semantic_guess(subtype, features)
            body_blocks = [
                {'region': _region_name(t['box']), 'box_pct': t['box']}
                for t in features['body_boxes']
            ]
            visual_blocks = [
                {'region': _region_name(v['box']), 'box_pct': v['box']}
                for v in features['large_visuals']
            ]
            styles_by_sig[signature] = {
                'id': f'content-style-{len(styles_by_sig) + 1:02d}',
                'subtype': subtype,
                'name': {
                    'content-text-image': '图文混排内容页',
                    'content-timeline': '时间轴/流程内容页',
                    'content-two-column-text': '多段正文内容页',
                    'content-multi-panel': '多模块内容页',
                    'content-general': '通用内容页',
                }.get(subtype, '内容页样式'),
                'slides': [],
                'description': _describe_content_style(subtype, features),
                'semantic_type_guess': semantic['semantic_type_guess'],
                'information_relation': semantic['information_relation'],
                'selection_rule': semantic['selection_rule'],
                'body_blocks': body_blocks,
                'visual_blocks': visual_blocks,
                'usage_rule': _usage_rule(subtype),
            }
        styles_by_sig[signature]['slides'].append(idx)

    return list(styles_by_sig.values())


# ─── 主入口 ──────────────────────────────────────────────────────────────────


def extract_structure(pptx_path: str, output_dir: Optional[str] = None) -> Dict[str, Any]:
    result: Dict[str, Any] = {
        'file': os.path.basename(pptx_path),
        'slide_count': 0,
        'slide_size': {},
        'colors': {},             # 主题色板（12色）
        'actual_colors': {},      # 实际使用色（含非主题色，按频次）
        'bg_text_mapping': {},    # 背景色 → 文字色映射
        'fonts': {},              # 主题字体（major/minor）
        'font_sizes': {},         # 各占位符层级实际字号 + 字体名
        'para_alignment': {},     # 各占位符层级段落对齐
        'master': {},             # 母版固定装饰元素
        'component_styles': {},   # 边框/阴影/圆角/渐变典型值
        'layouts': [],
        'slide_roles': [],        # 逐页角色：封面/目录/章节/内容/结束页
        'content_layout_styles': [],  # 内容页样式库（按结构去重）
    }

    if not os.path.exists(pptx_path):
        raise FileNotFoundError(f"PPTX 文件不存在: {pptx_path}")

    with zipfile.ZipFile(pptx_path, 'r') as zf:
        theme_files = [f for f in zf.namelist()
                       if f.startswith('ppt/theme/theme') and f.endswith('.xml')]
        if theme_files:
            theme_xml = zf.read(theme_files[0]).decode('utf-8')
            result['colors'] = extract_theme_colors(theme_xml)
            result['fonts']  = extract_theme_fonts(theme_xml)

        if 'ppt/presentation.xml' in zf.namelist():
            pres_xml = zf.read('ppt/presentation.xml').decode('utf-8')
            result['slide_size'] = extract_slide_size(pres_xml)

        slide_files = [f for f in zf.namelist()
                       if re.match(r'ppt/slides/slide\d+\.xml', f)]
        result['slide_count'] = len(slide_files)

    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        tc = result['colors']   # 主题色板，用于解析 schemeClr
        tf = result['fonts']    # 主题字体，用于解析 +mj-lt 等引用
        result['actual_colors'] = extract_actual_colors(prs, tc)
        result['bg_text_mapping'] = extract_bg_text_mapping(prs, tc)
        result['font_sizes'] = extract_font_sizes(prs, tf)
        result['para_alignment'] = extract_para_alignment(prs)
        result['master'] = extract_master_elements(prs, tc)
        result['component_styles'] = extract_component_styles(prs)
        result['layouts'] = extract_layouts(prs, tf)
        result['slide_roles'] = classify_slide_roles(prs)
        result['content_layout_styles'] = extract_content_layout_styles(prs, result['slide_roles'])
    except ImportError:
        print("警告: python-pptx 未安装，跳过深度提取", file=sys.stderr)
    except Exception as e:
        print(f"深度提取失败: {e}", file=sys.stderr)

    return result


def main():
    if len(sys.argv) < 2:
        print("用法: python extract_structure.py <pptx文件路径> [输出JSON路径]")
        sys.exit(1)

    pptx_path = sys.argv[1]
    output_path = sys.argv[2] if len(sys.argv) > 2 else None

    try:
        output_dir = os.path.dirname(os.path.abspath(output_path)) if output_path else None
        result = extract_structure(pptx_path, output_dir=output_dir)
        if output_dir:
            result['temp_dir'] = output_dir.replace('\\', '/')
        output_json = json.dumps(result, ensure_ascii=False, indent=2)
        if output_path:
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(output_json)
            print(f"结构数据已保存到: {output_path}")
        else:
            print(output_json)
    except Exception as e:
        print(f"错误: {e}", file=sys.stderr)
        sys.exit(1)


if __name__ == '__main__':
    main()
