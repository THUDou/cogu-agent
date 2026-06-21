import os
import sys
import pytest
from pptx import Presentation

# Allow importing from parent directory in test environment without package installation
sys.path.append(os.path.join(os.path.dirname(__file__), '..', 'scripts'))
from split_pptx import split_pptx


def make_test_pptx(path, num_slides=3):
    prs = Presentation()
    for _ in range(num_slides):
        prs.slides.add_slide(prs.slide_layouts[0])
    prs.save(path)


def test_split_creates_correct_files(tmp_path):
    input_path = str(tmp_path / "test.pptx")
    make_test_pptx(input_path, num_slides=3)
    output_dir = str(tmp_path / "output")

    count = split_pptx(input_path, output_dir)

    files = sorted(os.listdir(output_dir))
    assert files == ["slide_001.pptx", "slide_002.pptx", "slide_003.pptx"]
    assert count == 3


def test_split_each_file_has_one_slide(tmp_path):
    input_path = str(tmp_path / "test.pptx")
    make_test_pptx(input_path, num_slides=3)
    output_dir = str(tmp_path / "output")

    split_pptx(input_path, output_dir)

    for fname in ["slide_001.pptx", "slide_002.pptx", "slide_003.pptx"]:
        prs = Presentation(os.path.join(output_dir, fname))
        assert len(prs.slides) == 1


def test_split_max_slides_limits_output(tmp_path):
    input_path = str(tmp_path / "test.pptx")
    make_test_pptx(input_path, num_slides=5)
    output_dir = str(tmp_path / "output")

    count = split_pptx(input_path, output_dir, max_slides=2)

    files = sorted(os.listdir(output_dir))
    assert files == ["slide_001.pptx", "slide_002.pptx"]
    assert count == 2


def test_split_creates_output_dir_if_missing(tmp_path):
    input_path = str(tmp_path / "test.pptx")
    make_test_pptx(input_path, num_slides=1)
    output_dir = str(tmp_path / "new_dir" / "output")

    split_pptx(input_path, output_dir)

    assert os.path.isdir(output_dir)
    assert len(os.listdir(output_dir)) == 1
