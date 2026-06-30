import os
import sys

from pptx import Presentation
from pptx.util import Cm, Pt

sys.path.append(os.path.join(os.path.dirname(__file__), '..', 'scripts'))
from extract_structure import classify_slide_roles, extract_content_layout_styles


def add_textbox(slide, bounds, text, font_size=20):
    left, top, width, height = bounds
    box = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
    frame = box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    return box


def add_rect(slide, left, top, width, height):
    return slide.shapes.add_shape(1, Cm(left), Cm(top), Cm(width), Cm(height))


def build_role_fixture():
    prs = Presentation()
    blank = prs.slide_layouts[6]

    cover = prs.slides.add_slide(blank)
    add_textbox(cover, (4, 6, 18, 2), "学校介绍", 44)
    add_textbox(cover, (4, 8.5, 10, 1), "SCHOOL PROFILE", 18)

    toc = prs.slides.add_slide(blank)
    add_textbox(toc, (2, 1, 8, 1.2), "目录", 32)
    add_textbox(toc, (3, 4, 8, 1), "01 学校简介", 22)
    add_textbox(toc, (3, 6, 8, 1), "02 办学条件", 22)
    add_textbox(toc, (3, 8, 8, 1), "03 校园风采", 22)

    section = prs.slides.add_slide(blank)
    add_textbox(section, (4, 5, 4, 2), "01", 54)
    add_textbox(section, (8, 5.4, 9, 1.5), "学校简介", 34)

    content = prs.slides.add_slide(blank)
    add_textbox(content, (1.2, 0.7, 10, 1), "学校概况", 24)
    add_textbox(content, (1.8, 3, 11, 2.5), "第一段正文内容\n补充说明文字", 18)
    add_textbox(content, (1.8, 7, 11, 2.5), "第二段正文内容\n补充说明文字", 18)
    add_rect(content, 22, 2, 8, 5)

    closing = prs.slides.add_slide(blank)
    add_textbox(closing, (7, 6, 15, 2), "THANK YOU", 44)
    add_textbox(closing, (7, 8, 15, 1), "谢谢观看", 28)

    return prs


def test_classify_slide_roles_identifies_standard_page_types():
    prs = build_role_fixture()

    roles = classify_slide_roles(prs)

    assert [r["role"] for r in roles] == ["cover", "toc", "section", "content", "closing"]
    assert roles[3]["subtype"] == "content-text-image"
    assert roles[3]["confidence"] >= 0.6


def test_extract_content_layout_styles_deduplicates_similar_content_pages():
    prs = Presentation()
    blank = prs.slide_layouts[6]

    for _ in range(2):
        slide = prs.slides.add_slide(blank)
        add_textbox(slide, (1.2, 0.7, 10, 1), "页面标题", 24)
        add_textbox(slide, (1.8, 3, 11, 2.5), "第一段正文内容\n补充说明文字", 18)
        add_textbox(slide, (1.8, 7, 11, 2.5), "第二段正文内容\n补充说明文字", 18)
        add_rect(slide, 22, 2, 8, 5)

    slide = prs.slides.add_slide(blank)
    add_textbox(slide, (1.2, 0.7, 10, 1), "时间轴", 24)
    add_rect(slide, 3, 5, 24, 0.1)
    add_textbox(slide, (3, 6, 5, 1), "节点一", 16)
    add_textbox(slide, (12, 6, 5, 1), "节点二", 16)
    add_textbox(slide, (21, 6, 5, 1), "节点三", 16)

    roles = [
        {"page": 1, "role": "content", "subtype": "content-text-image"},
        {"page": 2, "role": "content", "subtype": "content-text-image"},
        {"page": 3, "role": "content", "subtype": "content-timeline"},
    ]

    styles = extract_content_layout_styles(prs, roles)

    assert len(styles) == 2
    assert styles[0]["slides"] == [1, 2]
    assert styles[0]["usage_rule"]
    assert styles[0]["semantic_type_guess"] == "image_text"
    assert styles[0]["information_relation"]
    assert styles[0]["selection_rule"]
    assert styles[1]["subtype"] == "content-timeline"
    assert styles[1]["semantic_type_guess"] in ("process", "timeline")


def test_content_word_does_not_trigger_toc_and_part_marker_is_section():
    prs = Presentation()
    blank = prs.slide_layouts[6]

    cover = prs.slides.add_slide(blank)
    add_textbox(cover, (2, 3, 16, 2), "Business Report", 38)
    add_textbox(cover, (2, 6, 20, 1), "your content is entered here", 16)
    add_textbox(cover, (2, 8, 10, 1), "汇报人：XXX", 14)

    toc = prs.slides.add_slide(blank)
    add_textbox(toc, (2, 2, 10, 1), "CONTENTS", 32)
    add_textbox(toc, (3, 5, 8, 1), "01 Overview", 20)

    section = prs.slides.add_slide(blank)
    add_textbox(section, (2, 3, 8, 1), "第一章节标题", 34)
    add_textbox(section, (2, 5, 12, 1), "工作内容汇报", 28)
    add_textbox(section, (2, 7, 18, 1), "your content is entered here", 14)

    roles = classify_slide_roles(prs)

    assert [r["role"] for r in roles] == ["cover", "toc", "section"]
