#!/usr/bin/env python3
"""Extract slide-level text, tables, notes, and basic shape metadata from PPTX."""

from __future__ import annotations

import argparse
import glob
import json
import logging
import os
from datetime import date, datetime
from pathlib import Path
from typing import Any, Dict, Iterable, List, Optional

from pptx import Presentation


EMU_PER_INCH = 914400
LOGGER = logging.getLogger(__name__)


def emu_to_in(value: Optional[int]) -> Optional[float]:
    if value is None:
        return None
    return round(value / EMU_PER_INCH, 3)


def shape_bounds(shape: Any) -> Dict[str, Optional[float]]:
    return {
        "left_in": emu_to_in(getattr(shape, "left", None)),
        "top_in": emu_to_in(getattr(shape, "top", None)),
        "width_in": emu_to_in(getattr(shape, "width", None)),
        "height_in": emu_to_in(getattr(shape, "height", None)),
    }


def placeholder_name(shape: Any) -> Optional[str]:
    try:
        return str(shape.placeholder_format.type)
    except Exception:
        return None


def iter_shapes(shapes: Iterable[Any]) -> Iterable[Any]:
    for shape in shapes:
        yield shape
        nested = getattr(shape, "shapes", None)
        if nested is not None:
            yield from iter_shapes(nested)


def clean_lines(text: str) -> List[str]:
    return [line.strip() for line in text.splitlines() if line.strip()]


def repair_utf8_mojibake(value: str) -> str:
    try:
        return value.encode("latin1").decode("utf-8")
    except UnicodeError:
        return value


def has_glob_pattern(value: str) -> bool:
    return any(char in value for char in "*?[")


def resolve_input_path(input_value: str) -> Path:
    candidates = [input_value]
    repaired = repair_utf8_mojibake(input_value)
    if repaired != input_value:
        candidates.append(repaired)

    for candidate in candidates:
        path = Path(candidate)
        if path.exists():
            return path

        if has_glob_pattern(candidate):
            matches = [
                Path(match)
                for match in glob.glob(candidate)
                if Path(match).suffix.lower() == ".pptx"
            ]
            if len(matches) == 1:
                return matches[0]
            if len(matches) > 1:
                rendered = "\n".join(str(match) for match in sorted(matches))
                raise ValueError(
                    "PPTX pattern matched multiple files; pass one explicit path:\n"
                    f"{rendered}"
                )

    raise FileNotFoundError(f"PPTX file not found: {input_value}")


def text_value_to_lines(value: Any) -> List[str]:
    if isinstance(value, str):
        return clean_lines(value)
    if isinstance(value, list):
        if all(isinstance(item, str) and len(item) <= 1 for item in value):
            return clean_lines("".join(value))
        return [
            line
            for item in value
            for line in clean_lines(str(item))
        ]
    if value is None:
        return []
    return clean_lines(str(value))


def json_safe_property(value: Any) -> Any:
    if isinstance(value, (date, datetime)):
        return value.isoformat()
    if value is None:
        return ""
    return value


def extract_document_properties(presentation: Presentation) -> Dict[str, Any]:
    props = presentation.core_properties
    fields = [
        "title",
        "author",
        "subject",
        "category",
        "keywords",
        "comments",
        "content_status",
        "content_type",
        "identifier",
        "language",
        "last_modified_by",
        "revision",
        "version",
        "created",
        "modified",
        "last_printed",
    ]
    return {
        field: json_safe_property(getattr(props, field, ""))
        for field in fields
    }


def extract_text_shape(shape: Any) -> Optional[Dict[str, Any]]:
    if not getattr(shape, "has_text_frame", False):
        return None

    paragraphs = []
    for paragraph in shape.text_frame.paragraphs:
        runs = [run.text for run in paragraph.runs if run.text]
        text = "".join(runs).strip() or paragraph.text.strip()
        if not text:
            continue
        paragraphs.append({
            "text": text,
            "level": paragraph.level,
        })

    if not paragraphs:
        return None

    return {
        "shape_id": getattr(shape, "shape_id", None),
        "name": getattr(shape, "name", ""),
        "placeholder": placeholder_name(shape),
        "bounds": shape_bounds(shape),
        "paragraphs": paragraphs,
        "text": "\n".join(item["text"] for item in paragraphs),
    }


def extract_table(shape: Any) -> Optional[Dict[str, Any]]:
    if not getattr(shape, "has_table", False):
        return None

    rows = []
    for row in shape.table.rows:
        rows.append([cell.text.strip() for cell in row.cells])

    return {
        "shape_id": getattr(shape, "shape_id", None),
        "name": getattr(shape, "name", ""),
        "bounds": shape_bounds(shape),
        "rows": rows,
    }


def extract_notes(slide: Any) -> str:
    try:
        notes = slide.notes_slide.notes_text_frame.text
    except Exception:
        return ""
    return "\n".join(clean_lines(notes))


def summarize_slide(slide: Any, slide_number: int) -> Dict[str, Any]:
    text_shapes = []
    tables = []
    image_count = 0
    chart_count = 0
    media_shape_names = []

    for shape in iter_shapes(slide.shapes):
        text_shape = extract_text_shape(shape)
        if text_shape:
            text_shapes.append(text_shape)

        table = extract_table(shape)
        if table:
            tables.append(table)

        if getattr(shape, "has_chart", False):
            chart_count += 1

        shape_type = str(getattr(shape, "shape_type", "")).upper()
        if "PICTURE" in shape_type:
            image_count += 1
            media_shape_names.append(getattr(shape, "name", ""))

    full_text_parts = []
    for text_shape in text_shapes:
        full_text_parts.extend(clean_lines(text_shape["text"]))
    for table in tables:
        for row in table["rows"]:
            row_text = " | ".join(cell for cell in row if cell)
            if row_text:
                full_text_parts.append(row_text)
    notes = extract_notes(slide)
    if notes:
        full_text_parts.append(f"Speaker notes: {notes}")

    title = ""
    for text_shape in text_shapes:
        placeholder = (text_shape.get("placeholder") or "").upper()
        if "TITLE" in placeholder and text_shape["paragraphs"]:
            title = text_shape["paragraphs"][0]["text"]
            break
    if not title and text_shapes:
        title = text_shapes[0]["paragraphs"][0]["text"]

    return {
        "slide_number": slide_number,
        "title_guess": title,
        "text_shapes": text_shapes,
        "tables": tables,
        "speaker_notes": notes,
        "image_count": image_count,
        "chart_count": chart_count,
        "media_shape_names": [name for name in media_shape_names if name],
        "full_text": "\n".join(full_text_parts),
    }


def extract_pptx(input_path: str) -> Dict[str, Any]:
    prs = Presentation(input_path)
    return {
        "schema_version": "pptx-content-summarizer-text-v2",
        "input_pptx": os.path.abspath(input_path),
        "document_properties": extract_document_properties(prs),
        "slide_count": len(prs.slides),
        "slide_size": {
            "width_in": emu_to_in(prs.slide_width),
            "height_in": emu_to_in(prs.slide_height),
        },
        "slides": [
            summarize_slide(slide, index + 1)
            for index, slide in enumerate(prs.slides)
        ],
    }


def compact_for_summary(data: Dict[str, Any]) -> Dict[str, Any]:
    """Return the small payload that downstream summary writing should read."""
    slides = []
    for slide in data.get("slides", []):
        content = "\n".join(text_value_to_lines(slide.get("full_text", "")))
        if not content:
            content_parts = []
            for text_shape in slide.get("text_shapes", []):
                shape_lines = text_value_to_lines(text_shape.get("text", ""))
                if not shape_lines:
                    for paragraph in text_shape.get("paragraphs", []):
                        shape_lines.extend(text_value_to_lines(paragraph.get("text", "")))
                content_parts.extend(shape_lines)
            for table in slide.get("tables", []):
                for row in table.get("rows", []):
                    row_text = " | ".join(str(cell) for cell in row if str(cell).strip())
                    if row_text:
                        content_parts.append(row_text)
            speaker_notes = "\n".join(text_value_to_lines(slide.get("speaker_notes", "")))
            if speaker_notes:
                content_parts.append(f"Speaker notes: {speaker_notes}")
            content = "\n".join(content_parts)

        slides.append({
            "slide_number": slide.get("slide_number"),
            "title": slide.get("title_guess", ""),
            "content": content,
            "tables": [
                table.get("rows", [])
                for table in slide.get("tables", [])
            ],
            "speaker_notes": slide.get("speaker_notes", ""),
            "image_count": slide.get("image_count", 0),
            "chart_count": slide.get("chart_count", 0),
        })

    return {
        "schema_version": "pptx-content-summarizer-summary-input-v2",
        "input_pptx": data.get("input_pptx", ""),
        "document_properties": data.get("document_properties", {}),
        "slide_count": data.get("slide_count", len(slides)),
        "slides": slides,
    }


def write_slide_inputs(summary_data: Dict[str, Any], output_dir: Path) -> None:
    """Write one input JSON file per slide for low-context-pressure page card generation."""
    output_dir.mkdir(parents=True, exist_ok=True)
    doc_props = summary_data.get("document_properties", {})
    slide_count = summary_data.get("slide_count", 0)
    slides = summary_data.get("slides", [])
    for slide in slides:
        slide_number = slide["slide_number"]
        filename = f"slide_{slide_number:03d}_input.json"
        payload = {
            "schema_version": "pptx-content-summarizer-slide-input-v1",
            "document_properties": doc_props,
            "slide_count": slide_count,
            "slide": slide,
        }
        (output_dir / filename).write_text(
            json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8"
        )
    LOGGER.info("Wrote %d slide input files to: %s", len(slides), output_dir.resolve())


def main() -> None:
    logging.basicConfig(level=logging.INFO, format="%(levelname)s:%(name)s:%(message)s")

    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("pptx", help="Input .pptx file")
    parser.add_argument("--output", "-o", required=True, help="Output detailed JSON path")
    parser.add_argument(
        "--slide-inputs",
        help="Directory to write per-slide input JSON files (slide_001_input.json, ...)",
    )
    args = parser.parse_args()

    input_path = resolve_input_path(args.pptx)
    if input_path.suffix.lower() != ".pptx":
        raise ValueError(f"Only .pptx files are supported: {input_path}")

    data = extract_pptx(str(input_path))
    payload = json.dumps(data, ensure_ascii=False, indent=2)

    output_path = Path(args.output)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_text(payload, encoding="utf-8")
    LOGGER.info("Wrote detailed extraction JSON: %s", output_path.resolve())

    if args.slide_inputs:
        write_slide_inputs(compact_for_summary(data), Path(args.slide_inputs))


if __name__ == "__main__":
    main()
