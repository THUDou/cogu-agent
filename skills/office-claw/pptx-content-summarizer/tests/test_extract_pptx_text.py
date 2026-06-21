from __future__ import annotations

import importlib.util
from datetime import datetime
from pathlib import Path

from pptx import Presentation


SCRIPT_PATH = (
    Path(__file__).resolve().parents[1]
    / "scripts"
    / "extract_pptx_text.py"
)


def load_module():
    spec = importlib.util.spec_from_file_location("extract_pptx_text", SCRIPT_PATH)
    module = importlib.util.module_from_spec(spec)
    assert spec.loader is not None
    spec.loader.exec_module(module)
    return module


def test_extract_pptx_includes_core_document_properties(tmp_path):
    module = load_module()

    pptx_path = tmp_path / "sample.pptx"
    prs = Presentation()
    props = prs.core_properties
    props.title = "Cloud Development Roadmap"
    props.author = "Strategy Office"
    props.subject = "Cloud growth analysis"
    props.category = "Executive report"
    props.keywords = "cloud;growth;roadmap"
    props.comments = "For leadership review"
    props.last_modified_by = "Analyst"
    props.revision = 7
    props.created = datetime(2026, 5, 24, 17, 6, 29)
    props.modified = datetime(2026, 5, 26, 9, 30, 0)
    prs.save(pptx_path)

    data = module.extract_pptx(str(pptx_path))

    assert data["document_properties"]["title"] == "Cloud Development Roadmap"
    assert data["document_properties"]["author"] == "Strategy Office"
    assert data["document_properties"]["subject"] == "Cloud growth analysis"
    assert data["document_properties"]["category"] == "Executive report"
    assert data["document_properties"]["keywords"] == "cloud;growth;roadmap"
    assert data["document_properties"]["comments"] == "For leadership review"
    assert data["document_properties"]["last_modified_by"] == "Analyst"
    assert data["document_properties"]["revision"] == 7
    assert data["document_properties"]["created"].startswith("2026-05-24T17:06:29")
    assert data["document_properties"]["modified"].startswith("2026-05-26T09:30:00")


def test_compact_for_summary_uses_v2_schema_and_document_properties():
    module = load_module()

    summary = module.compact_for_summary({
        "document_properties": {
            "title": "Cloud Development Roadmap",
            "author": "Strategy Office",
        },
        "slide_count": 1,
        "slides": [
            {
                "slide_number": 1,
                "title_guess": "Opening",
                "full_text": "Key metric: 42%",
                "tables": [],
                "speaker_notes": "",
                "image_count": 0,
                "chart_count": 0,
            }
        ],
    })

    assert summary["schema_version"] == "pptx-content-summarizer-summary-input-v2"
    assert summary["document_properties"] == {
        "title": "Cloud Development Roadmap",
        "author": "Strategy Office",
    }
    assert summary["slides"][0]["content"] == "Key metric: 42%"


def test_resolve_input_path_accepts_unexpanded_glob_pattern(tmp_path):
    module = load_module()

    docs_dir = tmp_path / "test_docs"
    docs_dir.mkdir()
    pptx_path = docs_dir / "AI RPA 金融版.pptx"
    Presentation().save(pptx_path)

    resolved = module.resolve_input_path(str(docs_dir / "*AI RPA*.pptx"))

    assert resolved == pptx_path


def test_resolve_input_path_repairs_utf8_mojibake_path(tmp_path):
    module = load_module()

    pptx_path = tmp_path / "智能流程机器人.pptx"
    Presentation().save(pptx_path)
    mojibake_path = str(pptx_path).encode("utf-8").decode("latin1")

    resolved = module.resolve_input_path(mojibake_path)

    assert resolved == pptx_path


def test_compact_for_summary_falls_back_to_shape_text_when_full_text_is_empty():
    module = load_module()

    summary = module.compact_for_summary({
        "slide_count": 1,
        "slides": [
            {
                "slide_number": 1,
                "title_guess": "Fallback",
                "full_text": "",
                "text_shapes": [
                    {
                        "text": "",
                        "paragraphs": [
                            {"text": "Primary point", "level": 0},
                            {"text": "Second point", "level": 1},
                        ],
                    }
                ],
                "tables": [
                    {"rows": [["Metric", "Meaning"], ["42%", "Example"]]},
                ],
                "speaker_notes": "Speaker context",
                "image_count": 0,
                "chart_count": 0,
            }
        ],
    })

    assert summary["slides"][0]["content"] == (
        "Primary point\n"
        "Second point\n"
        "Metric | Meaning\n"
        "42% | Example\n"
        "Speaker notes: Speaker context"
    )
